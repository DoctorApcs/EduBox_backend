import glob
import easyocr
import cv2
from pptx import Presentation

reader = easyocr.Reader(['en'])

def iou(boxA, boxB):
	xA = max(boxA[0], boxB[0])
	yA = max(boxA[1], boxB[1])
	xB = min(boxA[2], boxB[2])
	yB = min(boxA[3], boxB[3])

	interArea = max(0, xB - xA + 1) * max(0, yB - yA + 1)
	boxAArea = (boxA[2] - boxA[0] + 1) * (boxA[3] - boxA[1] + 1)
	boxBArea = (boxB[2] - boxB[0] + 1) * (boxB[3] - boxB[1] + 1)

	iou = interArea / float(boxAArea + boxBArea - interArea)
	return iou

def estimate_font(prs: Presentation, prs_image_folder: str):
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    prs_images_path = sorted(glob.glob('./metadata/img/*'), key=lambda x: int(x.split('/')[-1].rsplit('.')[-2]))

    for slide, prs_image_path in zip(prs.slides, prs_images_path):
        detected = reader.readtext(prs_image_path)
    
        img = cv2.imread('2.jpg')
        height, width, _ = img.shape
        bbox_list = []
        for bbox in detected:
            ([[x1,y1],[x2,y2],[x3,y3],[x4,y4]]), text, conf = bbox
            if conf > 0.5:
                bbox_list.append([int((y3-y2)*42/173), text]) # (42/173 is magic number)
        
        font_size_dict = dict()
        
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text != "":
                shape_bbox = [shape.left/slide_width, shape.top/slide_height, (shape.left + shape.width)/slide_width, (shape.top + shape.height)/slide_height]
            for bbox in bbox_list:
                if shape.text.strip() == bbox[1]:
                    font_size = bbox[0] # extract font size of the shape here
                    if shape.text.strip() not in font_size_dict:
                        font_size_dict[shape.text.strip()] = font_size
                        
        return font_size_dict