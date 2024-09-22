from pptx import Presentation

def _object_rels(obj):
    rels = obj.rels

    # Change required for python-pptx 0.6.22
    check_rels_content = [k for k in rels]
    if isinstance(check_rels_content.pop(), str):
        return [v for k, v in rels.items()]
    else:
        return [k for k in rels]


def _exp_add_slide(ppt, slide_layout):
    """
    Function to handle slide creation in the Presentation, to avoid issues caused by default implementation.

    :param slide_layout:
    :return:
    """

    def generate_slide_partname(self):
        """Return |PackURI| instance containing next available slide partname."""
        from pptx.opc.packuri import PackURI

        sldIdLst = self._element.get_or_add_sldIdLst()

        existing_rels = [k.target_partname for k in _object_rels(self)]
        partname_str = "/ppt/slides/slide%d.xml" % (len(sldIdLst) + 1)

        while partname_str in existing_rels:
            import random
            import string

            random_part = ''.join(random.choice(string.ascii_letters) for i in range(2))
            partname_str = "/ppt/slides/slide%s%d.xml" % (random_part, len(sldIdLst) + 1)

        return PackURI(partname_str)

    def add_slide_part(self, slide_layout):
        """
        Return an (rId, slide) pair of a newly created blank slide that
        inherits appearance from *slide_layout*.
        """
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.parts.slide import SlidePart

        partname = generate_slide_partname(self)
        slide_layout_part = slide_layout.part
        slide_part = SlidePart.new(partname, self.package, slide_layout_part)
        rId = self.relate_to(slide_part, RT.SLIDE)
        return rId, slide_part.slide

    def add_slide_ppt(self, slide_layout):
        rId, slide = add_slide_part(self.part, slide_layout)
        slide.shapes.clone_layout_placeholders(slide_layout)
        self._sldIdLst.add_sldId(rId)
        return slide

    # slide_layout = self.get_master_slide_layout(slide_layout)
    return add_slide_ppt(ppt.slides, slide_layout)


def copy_shapes(source, dest):
    """
    Helper to copy shapes handling edge cases.

    :param source:
    :param dest:
    :return:
    """
    from pptx.shapes.group import GroupShape
    import copy

    # Copy all existing shapes
    for shape in source:
        if isinstance(shape, GroupShape):
            group = dest.shapes.add_group_shape()
            group.name = shape.name
            group.left = shape.left
            group.top = shape.top
            group.width = shape.width
            group.height = shape.height
            group.rotation = shape.rotation

            # Recursive copy of contents
            copy_shapes(shape.shapes, group)

            # Fix offset
            cur_el = group._element.xpath(".//p:grpSpPr")[0]
            ref_el = shape._element.xpath(".//p:grpSpPr")[0]
            parent = cur_el.getparent()
            parent.insert(
                parent.index(cur_el) + 1,
                copy.deepcopy(ref_el)
            )
            parent.remove(cur_el)

            result = group
        elif hasattr(shape, "image"):
            import io

            # Get image contents
            content = io.BytesIO(shape.image.blob)
            result = dest.shapes.add_picture(
                content, shape.left, shape.top, shape.width, shape.height
            )
            result.name = shape.name
            result.crop_left = shape.crop_left
            result.crop_right = shape.crop_right
            result.crop_top = shape.crop_top
            result.crop_bottom = shape.crop_bottom
        elif hasattr(shape, "has_chart") and shape.has_chart:
            from .charts import clone_chart
            result = clone_chart(shape, dest)
        else:
            import copy

            newel = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(newel, "p:extLst")
            result = dest.shapes[-1]


def duplicate_slide(ppt, slide_index: int):
    """
    Duplicate the slide with the given number in presentation.
    Adds the new slide by default at the end of the presentation.

    :param ppt:
    :param slide_index: Slide number
    :return:
    """
    source = ppt.slides[slide_index]

    dest = _exp_add_slide(ppt, source.slide_layout)

    # Remove all shapes from the default layout
    for shape in dest.shapes:
        remove_shape(shape)

    # Copy all existing shapes
    copy_shapes(source.shapes, dest)

    # Copy all existing shapes
    if source.has_notes_slide:
        txt = source.notes_slide.notes_text_frame.text
        dest.notes_slide.notes_text_frame.text = txt

    return dest

def remove_shape(shape):
    """
    Helper to remove a specific shape.

    :source: https://stackoverflow.com/questions/64700638/is-there-a-way-to-delete-a-shape-with-python-pptx

    :param shape:
    :return:
    """
    el = shape.element  # --- get reference to XML element for shape
    el.getparent().remove(el)  # --- remove that shape element from its tree
    
def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    xml_slides.remove(slides[index]) 

def replace_text(shape, new_text):
    paragraph = shape.text_frame.paragraphs[0]
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text

if __name__ == "__main__":
    input_pptx = 'slide_template.pptx'
    output_pptx = 'output.pptx'
    prs = Presentation(input_pptx)
    duplicate_slide(prs, 1)
    # move_slide(prs, len(prs.slides)-1, 2)
    delete_slide(prs, 0)
    prs.save(output_pptx)
