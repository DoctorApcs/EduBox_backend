import os
import json
import copy
import openai
from pptx import Presentation
from pptx.util import Pt
from dotenv import load_dotenv
from utils.utils import *
from pre_generate import generate_outline
from utils.openai_api import get_chat_response
from utils.slide import duplicate_slide, delete_slide, replace_text
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Pt
load_dotenv(override=True)
api_key = os.getenv("OPENAI_API_KEY")

def extract_text_boxes(prs):
    all_slides_layout = []
    for slide_index, slide in enumerate(prs.slides):
        slide_layout = {
            "slide_number": slide_index + 1, 
            "shapes": {},
            "layout_type": slide.slide_layout.name  # Add this line
        }
        # Process placeholders first
        for placeholder in slide.placeholders:
            if hasattr(placeholder, 'text_frame'):
                max_chars = sum(len(p.text) for p in placeholder.text_frame.paragraphs)
                if max_chars > 0:
                    slide_layout["shapes"][f"Placeholder_{placeholder.placeholder_format.idx}"] = {
                        "text": placeholder.text_frame.text,
                        "max_chars": max_chars
                    }
        
        # Process shapes, but only if they're not already covered by a placeholder
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                max_chars = sum(len(p.text) for p in shape.text_frame.paragraphs)
                if max_chars > 0 and shape.text_frame.text not in [s["text"] for s in slide_layout["shapes"].values()]:
                    slide_layout["shapes"][shape.name] = {
                        "text": shape.text_frame.text,
                        "max_chars": max_chars
                    }
        
        all_slides_layout.append(slide_layout)
        
        
    return all_slides_layout

def classify_text_type(text):
    text = text.lower().strip()
    if text == "title" or text == "subtitle":
        return text
    elif "short body text" in text:
        return "short_body_text"
    elif "body text" in text:
        return "body_text"
    elif "section" in text:
        return text
    elif "header" in text:
        return "header"
    elif "awesome word" in text or "keyword" in text:
        return "awesome_word"
    elif "section number" in text:
        return "section number"
    elif "Table of contents" in text:
        return "Table of contents"
    else:
        return "other"

def get_type_description(text_type):
    descriptions = {
        "title": "Short, catchy phrase summarizing the main point (5-10 words)",
        "body_text": "Detailed information or explanation about the topic (50-200 words)",
        "short_body_text": "Brief explanation or key point related to the topic (10-20 words)",
        "section": "Short descriptive text for the section of the topic (3-5 words)",
        "header": "Short descriptive text for the section of the topic (3-5 words)",
        "awesome_word": "Single word or short phrase highlighting a key concept (1-3 words)",
    }
    return descriptions.get(text_type, "")

def refine_slide_layout(all_slides_layout):
    refined_layout = []
    for slide in all_slides_layout:
        refined_slide = {
            "slide_number": slide["slide_number"],
            "layout_type": slide["layout_type"],  # Add this line
            "shapes": {}
        }
        for shape_name, shape_info in slide["shapes"].items():
            text_type = classify_text_type(shape_info["text"])
            refined_slide["shapes"][shape_name] = {
                "type": text_type if text_type != "other" else "other",
            }
        refined_layout.append(refined_slide)
    return refined_layout

def generate_content_with_gpt(topic: str, refined_layout: dict, outline: list):
    # Create a mapping of slide numbers to outline points
    outline_mapping = {str(i+1): point for i, point in enumerate(outline)}

    prompt = f"""
    Generate content for a presentation on '{topic}'. The presentation structure, outline, and layout are as follows:

    Outline:
    {json.dumps(outline_mapping, indent=2)}

    Refined Layout:
    {json.dumps(refined_layout, indent=2)}

    Guidelines:
    1. Generate content for EVERY shape in EVERY slide, including all Placeholders and Google Shapes.
    2. All content must relate directly to '{topic}' and STRICTLY follow the provided outline.
    3. Ensure that the content of each slide aligns PERFECTLY with the corresponding outline entry and layout type.
    4. Use technical terms and concepts specific to '{topic}'.
    5. Adhere STRICTLY to these content type and word limit guidelines:
       - 'title': EXACTLY 3-5 words
       - 'subtitle': EXACTLY 3-5 words
       - 'header': EXACTLY 2-3 words (use short, impactful words to avoid text overflow)
       - 'short_body_text': EXACTLY 5-10 words
       - 'body_text': EXACTLY 20-30 words
       - 'awesome_word': EXACTLY 1-2 words
       - 'section': EXACTLY 1-2 words (ensure very short, concise text to fit within the shape)
       - 'other': Keep original text or generate EXACTLY 1-2 words if empty
    6. You MUST fill each shape to its MAXIMUM word limit to ensure the layout is fully utilized.
    7. For each shape, start the content with its type in square brackets, e.g., [title] or [body_text].
    8. Use common abbreviations for programming terms to save space and avoid overflow in shapes.
    9. Maintain a clear, logical flow throughout the presentation, following the sequence in the outline.
    10. For slides with multiple shape types, ensure content is cohesive and builds on each point.
    11. Use bullet points for body text where appropriate to enhance readability.

    CRITICAL:
    - You MUST generate content for EVERY shape listed in the refined layout for each slide, including all Placeholders and Google Shapes.
    - The output JSON structure MUST EXACTLY match the structure of the refined layout.
    - STRICTLY adhere to the word limits for each content type, using the MAXIMUM number of words allowed.
    - Do not skip any shapes, even if they seem redundant or unnecessary.
    - Ensure that the content for each slide PRECISELY follows the corresponding outline point.
    - FILL EACH SHAPE TO ITS MAXIMUM CAPACITY while staying within the word limit.

    Return a JSON object where:
    - Keys are slide numbers (as strings, from "0" to "{len(refined_layout) - 1}")
    - Values are objects with shape names (e.g., "Placeholder_0", "Google Shape;617;p56") as keys and generated content as values
    - INCLUDE ALL shape names from the refined layout, even if the content is minimal
    
    Ensure the content adheres strictly to the type and word limits for each shape type, avoids text overflow, and maintains logical consistency with the layout type and outline. MAXIMIZE the use of available space in each shape.
    """

    response = get_chat_response(prompt=prompt, api_key=api_key, temperature=0.3)
    
    return json.loads(response)

def adjust_font_size_to_fit(shape, text, max_chars):
    target_shape = shape
    original_font_size = target_shape.text_frame.paragraphs[0].font.size
    # target_shape.text_frame.clear()
    # target_shape.text_frame.text = text
    replace_text(target_shape, text)
    
    
    if len(text) > max_chars:
        new_font_size = Pt(original_font_size.pt * 0.85)  # Reduce to 90% of original size
        # print("New font size:", new_font_size)
        # print("Font size:", original_font_size.pt)
        for paragraph in target_shape.text_frame.paragraphs:
            paragraph.font.size = new_font_size

def update_slides(prs, generated_content, all_slides_layout):
    for slide_index, slide in enumerate(prs.slides):
        print(f"Processing slide {slide_index + 1}")
        slide_content = generated_content.get(str(slide_index), {})
        
        original_layout = all_slides_layout[slide_index]['shapes']
        
        for shape_name, shape_info in original_layout.items():
            target_shape = None
            if shape_name.startswith("Placeholder_"):
                placeholder_index = int(shape_name.split("_")[1])
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == placeholder_index:
                        target_shape = placeholder
                        break
            else:
                for shape in slide.shapes:
                    if shape.name == shape_name:
                        target_shape = shape
                        break

            if target_shape and hasattr(target_shape, 'text_frame'):
                shape_type = classify_text_type(shape_info['text'])
                if shape_type == 'other':
                    print(f"Skipping 'other' shape '{shape_name}' on slide {slide_index + 1}")
                    continue
                if shape_name in slide_content:
                    new_text = slide_content[shape_name]
                    
                    # Extract the content from the generated text
                    content_parts = new_text.split(']', 1)
                    if len(content_parts) > 1:
                        content = content_parts[1].strip()
                    else:
                        content = new_text.strip()
                    
                    try:
                        print(f"Updating shape '{shape_name}' on slide {slide_index + 1} with text: {content[:50]}...")
                        
                        # Set up the text frame
                        text_frame = target_shape.text_frame
                        text_frame.clear()  # Clear existing text
                        text_frame.word_wrap = True
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        
                        # Add the new text
                        p = text_frame.paragraphs[0]
                        p.text = content
                        
                        # Set font size based on shape type
                        # if shape_type in ['title', 'subtitle']:
                        #     font_size = Pt(24)
                        # elif shape_type in ['header', 'section']:
                        #     font_size = Pt(18)
                        # else:
                        #     font_size = Pt(14)
                        
                        # Apply font size to all runs in the paragraph
                        # for run in p.runs:
                        #     run.font.size = font_size
                        
                    except Exception as e:
                        print(f"Error updating shape '{shape_name}' on slide {slide_index + 1}: {str(e)}")
                else:
                    print(f"Warning: No content generated for shape '{shape_name}' on slide {slide_index + 1}")
            else:
                print(f"Warning: Shape '{shape_name}' not found or doesn't have a text frame on slide {slide_index + 1}")

    print("All slides processed.")

# def update_slides(prs, generated_content, all_slides_layout):
#     for slide_index, slide in enumerate(prs.slides):
#         print(f"Processing slide {slide_index + 1}")
#         slide_content = generated_content.get(str(slide_index), {})
        
#         original_layout = all_slides_layout[slide_index]['shapes']
        
#         for shape_name, shape_info in original_layout.items():
#             target_shape = None
#             if shape_name.startswith("Placeholder_"):
#                 placeholder_index = int(shape_name.split("_")[1])
#                 for placeholder in slide.placeholders:
#                     if placeholder.placeholder_format.idx == placeholder_index:
#                         target_shape = placeholder
#                         break
#             else:
#                 for shape in slide.shapes:
#                     if shape.name == shape_name:
#                         target_shape = shape
#                         break

#             if target_shape and hasattr(target_shape, 'text_frame'):
#                 target_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
#                 shape_type = classify_text_type(shape_info['text'])
#                 if shape_type == 'other':
#                     print(f"Skipping 'other' shape '{shape_name}' on slide {slide_index + 1}")
#                     continue
#                 if shape_name in slide_content:
#                     new_text = slide_content[shape_name]
                    
#                     # Extract the content from the generated text
#                     content_parts = new_text.split(']', 1)
#                     if len(content_parts) > 1:
#                         content = content_parts[1].strip()
#                     else:
#                         content = new_text.strip()
                    
#                     try:
#                         print(f"Updating shape '{shape_name}' on slide {slide_index + 1} with text: {content[:50]}...")
#                         target_shape.text_frame.text = content
#                     except Exception as e:
#                         print(f"Error updating shape '{shape_name}' on slide {slide_index + 1}: {str(e)}")
#                 else:
#                     print(f"Warning: No content generated for shape '{shape_name}' on slide {slide_index + 1}")
#             else:
#                 print(f"Warning: Shape '{shape_name}' not found or doesn't have a text frame on slide {slide_index + 1}")

#     print("All slides processed.")

def pptx_layout_design(topic: str, template_pptx_path: str, slides_layout: dict, prs: Presentation, num_slides: int = 10):
    # Slide layout to dict
    slides_layout_dict = {}
    for slide_layout in slides_layout:
        slides_layout_dict[f"slide_{int(slide_layout['slide_number'])-1}"] = copy.deepcopy(slide_layout)

    # Generate outline
    outline = generate_outline(topic=topic, template_pptx_path=template_pptx_path, num_slides=num_slides)
    #outline = load_json('./src/metadata/outline_test.json')
    
    
    outline_list = []
    new_slides_layout = []
    origin_prs_length = len(prs.slides)
    for slide_number, slide_content in enumerate(list(outline.values())):
        tmp = copy.deepcopy(slides_layout_dict[slide_content['template_slide_number']]) 
        tmp['slide_number'] = slide_number
        
        new_slides_layout.append(tmp)

        outline_list.append({
            'slide_number': slide_number + 1,
            'outline': slide_content['outline']
        })
        
        
        duplicate_slide(prs, int(slide_content['template_slide_number'].replace("slide_", "")))
    
    for _ in range(origin_prs_length):
        delete_slide(prs, 0)
    
    return new_slides_layout, outline_list, prs

def generate_presentation(prompt: str, template: str = None, save_path: str = './results/demo.pptx'):
    openai.api_key = api_key
    topic = prompt
    try:
        prs = Presentation(template)
        
        all_slides_layout = extract_text_boxes(prs)
        save_json(all_slides_layout, './src/metadata/slides_layout.json')
        all_slides_layout, outline_list, prs = pptx_layout_design(topic=topic, template_pptx_path=template, 
                                                                  slides_layout=all_slides_layout, prs=prs, num_slides=15)     
        
        save_json(all_slides_layout, './src/metadata/new_slides_layout.json')
        refined_layout = refine_slide_layout(all_slides_layout)
        save_json(refined_layout, './src/metadata/refined_slides_layout.json')
        print("Refined slides layout metadata saved to metadata/refined_slides_layout.json")
        
        generated_content = generate_content_with_gpt(topic, refined_layout, outline_list)
        save_json(generated_content, './src/metadata/generated_content.json')
        print("Generated content saved to metadata/generated_content.json")
        
        update_slides(prs, generated_content, all_slides_layout)
        
        prs.save(save_path)
        print(f"Presentation updated successfully and saved as {save_path}")
    except Exception as e:
        print(f"An error occurred: {str(e)}")

if __name__ == "__main__":
    user_prompt = "Introduction to Python"  # Example prompt, replace with actual user input
    template_path = "./template/slide_template_2.pptx"  # Example template path, replace if needed
    save_file_path = "./results/demo.pptx"  # Example save path, replace if needed
    
    generate_presentation(prompt=user_prompt, template=template_path, save_path=save_file_path)

