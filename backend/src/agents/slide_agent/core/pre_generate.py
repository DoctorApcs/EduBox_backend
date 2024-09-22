import os
import json
import glob
from tqdm import tqdm
from dotenv import load_dotenv
from utils.utils import *
from utils.openai_api import get_vision_response, get_chat_response
from pptx import Presentation
import json

load_dotenv(override=True)
api_key = os.getenv("OPENAI_API_KEY")
    
def get_slide_description(slide_path: str):
    prompt = f"""Please generate a short but concise description of the content for this slide. Additionally, specify the slide's purpose (e.g., table of contents, first slide, ending slide, etc.), so I can determine whether this layout fit for my content.
    The response MUST be in format of:
    ## Slide Description:
    Short description of the slide including: text and main image position.
    ## Slide Purpose:
    The slide's purpose (e.g., table of contents, slide to introduction section, first slide, ending slide, etc.).
    """
    response = get_vision_response(image_path=slide_path, prompt=prompt, api_key=api_key)
    return response
    
def get_pptx_description(pptx_image_folder: str, save_path: str):
    slide_paths = sorted(glob.glob(pptx_image_folder + '*.jpg'), key=lambda x: int(x.split('/')[-1].replace(".jpg", "")))

    pptx_description = {}
    for slide_number, slide_path in enumerate(tqdm(slide_paths)):
        slide_description = get_slide_description(slide_path=slide_path)
        pptx_description[f'slide_{slide_number}'] = slide_description
        save_json(pptx_description, save_path)

def extract_slide_info(pptx_path):
    prs = Presentation(pptx_path)
    slides_info = []
    total_slides = len(prs.slides)
    
    for idx, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": idx,
            "shapes": {},
            "layout_type": slide.slide_layout.name,
            "text_content": "",
            "slide_type": ""
        }
        
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = shape.text.strip()
                if text:
                    slide_info["text_content"] += text + " "
                    slide_info["shapes"][shape.name] = {
                        "text": text,
                        "type": str(shape.shape_type)
                    }
        
        # Classify the slide and add the classification to slide_info
        slide_info["slide_type"] = classify_slide(slide_info, total_slides)
        
        slides_info.append(slide_info)
    
    return slides_info

def classify_slide(slide_info, total_slides):
    text_content = slide_info["text_content"].lower()
    
    if slide_info["slide_number"] == 0:
        return "title_slide"
    if slide_info["slide_number"] == total_slides - 1:
        return "end_slide"
    
    if "table of contents" in text_content:
        return "table_of_contents"
    
    if "section" in text_content:
        return "section_intro"
    
    if len([shape for shape in slide_info["shapes"].values() if "body text" in shape["text"].lower()]) >= 2:
        return "comparison_slide"
    
    if any("short body text" in shape["text"].lower() for shape in slide_info["shapes"].values()):
        return "bullet_point_slide"
    
    if "image" in text_content:
        return "image_with_caption"
    
    if "body text" in text_content:
        return "content_slide"
    
    if any(keyword in text_content for keyword in ["thank", "questions", "q&a"]):
        return "end_slide"
    
    return "other"

def generate_outline(topic: str, template_pptx_path: str, num_slides: int):
    template_description = extract_slide_info(template_pptx_path)
        
    slide_type_mapping = {}
    for slide in template_description:
        if slide['slide_type'] not in slide_type_mapping:
            slide_type_mapping[slide['slide_type']] = []
        slide_type_mapping[slide['slide_type']].append(f"slide_{slide['slide_number']}")
    
    # Find the last slide number
    last_slide_number = max(int(slide['slide_number']) for slide in template_description)
    
    # Add descriptions for each slide type
    slide_type_descriptions = {
        "title_slide": "The opening slide of the presentation, featuring the main topic.",
        "table_of_contents": "An overview slide listing the main sections of the presentation.",
        "section_intro": "A slide introducing a new main section or topic in the presentation.",
        "content_slide": "A slide containing detailed information, usually with paragraphs or long-form text.",
        "bullet_point_slide": "A slide with concise points or items, typically in a bulleted or numbered list.",
        "comparison_slide": "A slide that contrasts two or more ideas, concepts, or data points.",
        "image_with_caption": "A slide featuring a prominent image or visual element with an explanatory caption.",
        "end_slide": "The final slide of the presentation, summarizing key points or providing a conclusion."
    }
    
    prompt = f"""Generate a detailed and logically structured outline for a PowerPoint presentation on this topic:
    "{topic}"
    
    The presentation MUST have EXACTLY {num_slides} slides. No more, no less.
    
    Available slide types and their corresponding slide numbers:
    {json.dumps(slide_type_mapping, indent=2)}
    
    Slide type descriptions:
    {json.dumps(slide_type_descriptions, indent=2)}
    
    Detailed slide information:
    {json.dumps(template_description, indent=2)}
    
    Please follow these guidelines:
    1. Start with a title slide (slide_type: title_slide) that clearly states the topic.
    2. Include a table of contents slide (slide_type: table_of_contents) that outlines the main sections of the presentation.
    3. The outline MUST follow the structure defined in the table of contents. Each main section from the table of contents should be represented in the outline.
    4. For each main section:
       a. Begin with a section introduction slide (slide_type: section_intro).
       b. Use content slides (slide_type: content_slide) for detailed information.
       c. Use bullet point slides (slide_type: bullet_point_slide) for lists or brief points.
       d. Use comparison slides (slide_type: comparison_slide) to contrast different ideas or approaches when appropriate.
       e. Use image with caption slides (slide_type: image_with_caption) for visual elements that support the content.
    5. Ensure a logical progression within each section, moving from general concepts to specific details.
    6. End with a conclusion slide summarizing key points from all sections.
    7. The final slide should be an end_slide type, but focus its content on summarizing key takeaways rather than just saying "Thank you" or "Q&A".

    Ensure a logical flow between slides and sections, with clear transitions. The outline should reflect a coherent narrative that follows the table of contents structure.

    IMPORTANT: 
    - You MUST choose a template_slide_number that exists in the provided slide type mapping. Do not use any slide numbers that are not listed for each slide type.
    - Use the early slide numbers (e.g., slide_1, slide_2, slide_3) for the opening slides of your presentation.
    - The last slide of your outline MUST use "slide_{last_slide_number}" as its template_slide_number.
    - Ensure that the outline reflects the logical structure presented in the table of contents.
    - Each outline entry should be a brief, clear description of the slide's main content or purpose.
    - Focus on substantive content throughout the presentation, including the final slide.
    
    The response MUST be in the following JSON format:
    {{
        "1": {{
            "outline": "Brief description of slide content",
            "template_slide_number": "slide_X"
        }},
        ...
        "{num_slides}": {{
            "outline": "Brief description of slide content",
            "template_slide_number": "slide_{last_slide_number}"
        }}
    }}
    
    Where "template_slide_number" is one of the slide numbers provided that matches the intended content and layout type. Ensure that you use each slide type appropriately and don't overuse any particular type.
    
    Remember to maintain consistency between the table of contents and the actual content of the slides. Each main point in the table of contents should correspond to a section in the presentation.
    """

    response = get_chat_response(prompt=prompt, api_key=api_key)
    outline = json.loads(response)
    
    save_json(outline, 'src/metadata/outline_test.json')
    
    return outline

def generate_outline_with_binning(topic: str, template_pptx_path: str, bin_type: str):
    if bin_type == "0-10":
        num_slides = 10
    elif bin_type == "11-20":
        num_slides = 20
    elif bin_type == "21-30":
        num_slides = 30
    else:  # over 30
        num_slides = 40
    
    return generate_outline(topic, template_pptx_path, num_slides)

if __name__ == "__main__":
    generate_outline(topic="Introduction to Python Programming", template_pptx_path='/Users/mf/Desktop/Text2PPT/template/slide_template_2.pptx', num_slides=15)
