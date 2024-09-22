import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
import json
import uuid


class LayoutAnalysisAgent:
    def __init__(self, topic):
        self.topic = topic

    def extract_slide_info(self, presentation_path, template_name):
        prs = Presentation(presentation_path)
        slide_info = []

        for slide_number, slide in enumerate(prs.slides):
            print(f"Processing slide {slide_number + 1}...")
            shapes_info = []

            master = slide.slide_layout.slide_master
            layout = slide.slide_layout
            
            for shape in slide.shapes:
                shape_details = self.extract_shape_details(shape, prs.slide_width, prs.slide_height, master, layout)
                shapes_info.append(shape_details)

            slide_metadata = {
                "slide_dimensions": {
                    "width": prs.slide_width,
                    "height": prs.slide_height
                },
                "shapes": shapes_info
            }

            self.save_layout_to_json(slide_metadata, self.topic, template_name, slide_number + 1)
            slide_info.append(slide_metadata)

        return slide_info

    def extract_shape_details(self, shape, slide_width, slide_height, master, layout):
        details = {
            "id": str(shape._element.get("id") or uuid.uuid4()),
            "type": self.classify_shape_type(shape),
            "content_type": self.classify_shape_content_type(shape),
            "position": {
                "left": shape.left,
                "top": shape.top
            },
            "size": {
                "width": shape.width,
                "height": shape.height
            },
            "description": self.generate_shape_description(shape, slide_width, slide_height)
        }

        if shape.has_text_frame or shape.is_placeholder:
            details['text_properties'] = self.extract_text_properties(shape, master, layout)

        return details

    def classify_shape_type(self, shape):
        shape_types = {
            MSO_SHAPE_TYPE.AUTO_SHAPE: 'auto_shape',
            MSO_SHAPE_TYPE.CALLOUT: 'callout',
            MSO_SHAPE_TYPE.CHART: 'chart',
            MSO_SHAPE_TYPE.COMMENT: 'comment',
            MSO_SHAPE_TYPE.DIAGRAM: 'diagram',
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: 'embedded_ole_object',
            MSO_SHAPE_TYPE.GROUP: 'group',
            MSO_SHAPE_TYPE.IGX_GRAPHIC: 'igx_graphic',
            MSO_SHAPE_TYPE.LINE: 'line',
            MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: 'linked_ole_object',
            MSO_SHAPE_TYPE.LINKED_PICTURE: 'linked_picture',
            MSO_SHAPE_TYPE.MEDIA: 'media',
            MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT: 'ole_control_object',
            MSO_SHAPE_TYPE.PICTURE: 'picture',
            MSO_SHAPE_TYPE.PLACEHOLDER: 'placeholder',
            MSO_SHAPE_TYPE.SCRIPT_ANCHOR: 'script_anchor',
            MSO_SHAPE_TYPE.TABLE: 'table',
            MSO_SHAPE_TYPE.TEXT_BOX: 'text_box',
            MSO_SHAPE_TYPE.WEB_VIDEO: 'web_video'
        }
        return shape_types.get(shape.shape_type, 'other')

    def classify_shape_content_type(self, shape):
        if shape.has_text_frame:
            return self.classify_text_content(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        else:
            return "other"

    def classify_text_content(self, shape):
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type == PP_PLACEHOLDER.TITLE:
                return "title"
            elif placeholder_type == PP_PLACEHOLDER.BODY:
                return "body"
            elif placeholder_type == PP_PLACEHOLDER.CENTER_TITLE:
                return "center_title"
            elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                return "subtitle"
            else:
                return "other_text"
        else:
            return "text_box"

    def extract_text_properties(self, shape, master, layout):
        properties = {}

        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            master_placeholder = self.find_placeholder(master, ph_type)
            layout_placeholder = self.find_placeholder(layout, ph_type)

            properties.update(self.get_placeholder_properties(master_placeholder))
            properties.update(self.get_placeholder_properties(layout_placeholder))

        if shape.has_text_frame:
            text_frame = shape.text_frame
            if text_frame.paragraphs:
                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else None
                
                if first_run:
                    properties.update({
                        "font_name": first_run.font.name,
                        "font_size": first_run.font.size.pt if first_run.font.size else None,
                        "color": self.get_color_safely(first_run.font.color),
                        "bold": first_run.font.bold,
                        "italic": first_run.font.italic,
                        "underline": first_run.font.underline
                    })
                
                properties["alignment"] = str(first_paragraph.alignment) if first_paragraph.alignment else "LEFT"
                properties["max_chars"] = sum(len(p.text) for p in text_frame.paragraphs)

        return properties

    def find_placeholder(self, slide, placeholder_type):
        for shape in slide.placeholders:
            if shape.placeholder_format.type == placeholder_type:
                return shape
        return None

    def get_placeholder_properties(self, placeholder):
        if not placeholder or not placeholder.has_text_frame:
            return {}

        text_frame = placeholder.text_frame
        if not text_frame.paragraphs:
            return {}

        first_paragraph = text_frame.paragraphs[0]
        first_run = first_paragraph.runs[0] if first_paragraph.runs else None
        
        properties = {}
        if first_run:
            properties = {
                "font_name": first_run.font.name,
                "font_size": first_run.font.size.pt if first_run.font.size else None,
                "color": self.get_color_safely(first_run.font.color),
                "bold": first_run.font.bold,
                "italic": first_run.font.italic,
                "underline": first_run.font.underline
            }
        
        properties["alignment"] = str(first_paragraph.alignment) if first_paragraph.alignment else "LEFT"
        return properties

    def generate_shape_description(self, shape, slide_width, slide_height):
        content_type = self.classify_shape_content_type(shape)
        if content_type == "image":
            return self.describe_picture(shape, slide_width, slide_height)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return f"Group containing {len(shape.shapes)} elements"
        elif content_type in ["chart", "table"]:
            return f"{content_type.capitalize()} element"
        elif shape.has_text_frame:
            return self.describe_text_box(shape, slide_width, slide_height)
        else:
            return "Decorative shape or element"

    def describe_picture(self, shape, slide_width, slide_height):
        relative_width = shape.width / slide_width
        relative_height = shape.height / slide_height
        position = "center"
        if shape.left < slide_width * 0.25:
            position = "left"
        elif shape.left > slide_width * 0.75:
            position = "right"
        
        if relative_width > 0.5 and relative_height > 0.5:
            size = "large"
        elif relative_width < 0.2 or relative_height < 0.2:
            size = "small"
        else:
            size = "medium"
        
        return f"{size.capitalize()} image on the {position} of the slide"

    def describe_text_box(self, shape, slide_width, slide_height):
        content_type = self.classify_text_content(shape)
        relative_top = shape.top / slide_height
        relative_width = shape.width / slide_width

        if content_type in ["title", "center_title"]:
            return "Slide title"
        elif content_type == "subtitle":
            return "Slide subtitle"
        elif content_type == "body":
            return "Main body text"
        elif relative_top < 0.2:
            return "Header text"
        elif relative_top > 0.8:
            return "Footer text"
        elif relative_width > 0.5:
            return "Main content text"
        else:
            return "Supplementary text or caption"

    def extract_text_properties(self, shape, master, layout):
        properties = {}

        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            master_placeholder = self.find_placeholder(master, ph_type)
            layout_placeholder = self.find_placeholder(layout, ph_type)

            # Try to get properties from layout first, then master
            properties.update(self.get_placeholder_properties(layout_placeholder))
            if not all(properties.values()):
                properties.update(self.get_placeholder_properties(master_placeholder))

        if shape.has_text_frame:
            text_frame = shape.text_frame
            if text_frame.paragraphs:
                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else None
                
                if first_run:
                    properties.update({
                        "font_name": first_run.font.name or properties.get("font_name"),
                        "font_size": first_run.font.size.pt if first_run.font.size else properties.get("font_size"),
                        "color": self.get_color_safely(first_run.font.color) or properties.get("color"),
                        "bold": first_run.font.bold if first_run.font.bold is not None else properties.get("bold"),
                        "italic": first_run.font.italic if first_run.font.italic is not None else properties.get("italic"),
                        "underline": first_run.font.underline if first_run.font.underline is not None else properties.get("underline")
                    })
                
                properties["alignment"] = str(first_paragraph.alignment) if first_paragraph.alignment else properties.get("alignment", "LEFT")
                properties["max_chars"] = sum(len(p.text) for p in text_frame.paragraphs)

        # If properties are still None, try to get default values
        properties["font_name"] = properties["font_name"] or self.get_default_font_name(master)
        properties["font_size"] = properties["font_size"] or self.get_default_font_size(master)
        properties["color"] = properties["color"] or self.get_default_font_color(master)

        return properties

    def get_placeholder_properties(self, placeholder):
        if not placeholder or not placeholder.has_text_frame:
            return {}

        text_frame = placeholder.text_frame
        if not text_frame.paragraphs:
            return {}

        first_paragraph = text_frame.paragraphs[0]
        first_run = first_paragraph.runs[0] if first_paragraph.runs else None
        
        properties = {}
        if first_run:
            properties = {
                "font_name": first_run.font.name,
                "font_size": first_run.font.size.pt if first_run.font.size else None,
                "color": self.get_color_safely(first_run.font.color),
                "bold": first_run.font.bold,
                "italic": first_run.font.italic,
                "underline": first_run.font.underline
            }
        
        properties["alignment"] = str(first_paragraph.alignment) if first_paragraph.alignment else "LEFT"
        return properties

    def get_default_font_name(self, master):
        try:
            return master.slide_master.placeholders[0].text_frame.paragraphs[0].runs[0].font.name
        except:
            return "Arial"  # Default fallback

    def get_default_font_size(self, master):
        try:
            size = master.slide_master.placeholders[0].text_frame.paragraphs[0].runs[0].font.size
            return size.pt if size else 18  # Default to 18 if size is not set
        except:
            return 18  # Default fallback

    def get_default_font_color(self, master):
        try:
            color = master.slide_master.placeholders[0].text_frame.paragraphs[0].runs[0].font.color.rgb
            return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"
        except:
            return "#000000"  # Default to black

    def get_color_safely(self, color):
        if isinstance(color, RGBColor):
            return f"#{color.r:02x}{color.g:02x}{color.b:02x}"
        return None

    def save_layout_to_json(self, slide_metadata, topic, template_name, slide_number):
        folder_path = f"../layouts/{topic}/{template_name}"
        os.makedirs(folder_path, exist_ok=True)
        json_path = f"{folder_path}/layout_metadata_slide_{slide_number}.json"
        
        with open(json_path, 'w') as json_file:
            json.dump(slide_metadata, json_file, indent=4)
        print(f"Layout metadata saved to '{json_path}'.")

# Usage
topic = "Finance"
template_name = "template_11"
layout_agent = LayoutAnalysisAgent(topic)
layout_agent.extract_slide_info("../test_3.pptx", template_name)